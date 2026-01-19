"""
Attachment Processor Service

Extracts text content from various attachment types for indexing.
"""

import io
from pathlib import Path

from loguru import logger


class AttachmentProcessorError(Exception):
    """Base exception for attachment processing errors."""

    pass


class UnsupportedFormatError(AttachmentProcessorError):
    """Attachment format is not supported."""

    pass


class AttachmentProcessor:
    """
    Service for extracting text content from email attachments.

    Supports:
    - PDF files (.pdf)
    - Microsoft Word documents (.docx, .doc)
    - Microsoft Excel spreadsheets (.xlsx, .xls)
    - Microsoft PowerPoint presentations (.pptx, .ppt)
    - Plain text files (.txt, .csv, .json, .xml, .md, .yaml, etc.)
    - HTML files (.html, .htm)
    - ZIP archives (extracts text from supported files inside)

    Uses magic byte detection to identify file types when filenames
    or MIME types are incorrect/missing.
    """

    @staticmethod
    def sanitize_text_for_db(text: str) -> str:
        """
        Sanitize text for PostgreSQL storage.

        Removes null bytes and other characters that PostgreSQL
        cannot store in TEXT columns.

        Args:
            text: Raw text that may contain invalid characters

        Returns:
            Sanitized text safe for PostgreSQL storage
        """
        if not text:
            return text

        # Remove null bytes (PostgreSQL doesn't accept \x00)
        text = text.replace('\x00', '')

        # Remove other control characters except newlines, tabs, carriage returns
        # PostgreSQL can store these, but they can cause issues in some contexts
        import re
        # Keep \n (0x0A), \r (0x0D), \t (0x09), remove other control chars (0x00-0x08, 0x0B-0x0C, 0x0E-0x1F)
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)

        return text

    # Supported MIME types and their processors
    SUPPORTED_TYPES = {
        "application/pdf": "extract_pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "extract_docx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "extract_xlsx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "extract_pptx",
        "application/msword": "extract_doc",
        "application/vnd.ms-excel": "extract_xls",
        "application/vnd.ms-powerpoint": "extract_ppt",
        "text/plain": "extract_text",
        "text/html": "extract_html",
        "text/csv": "extract_text",
        "application/json": "extract_text",
        "application/xml": "extract_text",
        "text/xml": "extract_text",
        "application/zip": "extract_zip_contents",
    }

    # File extension to MIME type mapping
    EXTENSION_MAP = {
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".doc": "application/msword",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".xls": "application/vnd.ms-excel",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".ppt": "application/vnd.ms-powerpoint",
        ".txt": "text/plain",
        ".html": "text/html",
        ".htm": "text/html",
        ".csv": "text/csv",
        ".json": "application/json",
        ".xml": "application/xml",
        ".md": "text/plain",
        ".log": "text/plain",
        ".ini": "text/plain",
        ".cfg": "text/plain",
        ".yaml": "text/plain",
        ".yml": "text/plain",
        ".zip": "application/zip",
    }

    # Magic byte signatures for file type detection
    MAGIC_BYTES = {
        b"%PDF": "application/pdf",
        b"PK\x03\x04": "application/zip",  # ZIP-based (DOCX, XLSX, PPTX)
        b"\xd0\xcf\x11\xe0": "application/msword",  # Old Office format
    }

    def __init__(self, max_text_length: int = 100000):
        """
        Initialize attachment processor.

        Args:
            max_text_length: Maximum text length to extract (default 100K chars)
        """
        self.max_text_length = max_text_length

    def detect_mime_from_content(self, content: bytes) -> str | None:
        """
        Detect MIME type from file content using magic bytes.

        Args:
            content: File content bytes

        Returns:
            Detected MIME type or None
        """
        if len(content) < 4:
            return None

        header = content[:8]

        # Check magic bytes
        for magic, mime_type in self.MAGIC_BYTES.items():
            if header.startswith(magic):
                # For ZIP-based files, try to determine actual type
                if mime_type == "application/zip":
                    return self._detect_office_type(content)
                return mime_type

        return None

    def _detect_office_type(self, content: bytes) -> str:
        """
        Detect specific Office document type from ZIP content.

        Args:
            content: ZIP file content

        Returns:
            Specific MIME type for Office documents
        """
        import zipfile
        import io

        try:
            with zipfile.ZipFile(io.BytesIO(content), 'r') as zf:
                names = zf.namelist()
                if any('word/' in n for n in names):
                    return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                elif any('xl/' in n for n in names):
                    return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                elif any('ppt/' in n for n in names):
                    return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        except Exception:
            pass

        return "application/zip"

    def can_process(self, filename: str, content_type: str | None = None, content: bytes | None = None) -> bool:
        """
        Check if the attachment can be processed.

        Args:
            filename: Attachment filename
            content_type: MIME type (optional)
            content: File content for magic byte detection (optional)

        Returns:
            True if the attachment can be processed
        """
        # Check by content type (skip generic types)
        if content_type and content_type != "application/octet-stream" and content_type in self.SUPPORTED_TYPES:
            return True

        # Check by extension
        ext = Path(filename).suffix.lower()
        if ext in self.EXTENSION_MAP:
            return True

        # Try magic byte detection if content is provided
        if content:
            detected_type = self.detect_mime_from_content(content)
            if detected_type and detected_type in self.SUPPORTED_TYPES:
                return True

        return False

    def extract_text_from_attachment(
        self,
        content: bytes,
        filename: str,
        content_type: str | None = None,
    ) -> str:
        """
        Extract text from an attachment.

        Args:
            content: Attachment binary content
            filename: Attachment filename
            content_type: MIME type (optional)

        Returns:
            Extracted text content

        Raises:
            UnsupportedFormatError: If format is not supported
            AttachmentProcessorError: If extraction fails
        """
        # Determine the processor to use
        processor_name = None

        # Try content type first (skip generic types)
        if content_type and content_type != "application/octet-stream":
            processor_name = self.SUPPORTED_TYPES.get(content_type)

        # Fall back to extension
        if not processor_name:
            ext = Path(filename).suffix.lower()
            mime_type = self.EXTENSION_MAP.get(ext)
            if mime_type:
                processor_name = self.SUPPORTED_TYPES.get(mime_type)

        # Fall back to magic byte detection
        if not processor_name:
            detected_type = self.detect_mime_from_content(content)
            if detected_type:
                processor_name = self.SUPPORTED_TYPES.get(detected_type)
                logger.info(f"Detected {filename} as {detected_type} via magic bytes")

        if not processor_name:
            raise UnsupportedFormatError(
                f"Unsupported attachment format: {filename} ({content_type})"
            )

        # Get the processor method
        processor = getattr(self, processor_name, None)
        if not processor:
            raise UnsupportedFormatError(f"No processor for: {processor_name}")

        try:
            text = processor(content)

            # Sanitize text for PostgreSQL storage (remove null bytes, etc.)
            text = self.sanitize_text_for_db(text)

            # Truncate if necessary
            if len(text) > self.max_text_length:
                text = text[: self.max_text_length] + "\n[Content truncated...]"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from {filename}: {e}")
            raise AttachmentProcessorError(f"Failed to extract text: {e}") from e

    def extract_pdf(self, content: bytes) -> str:
        """Extract text from PDF file."""
        try:
            import pypdf

            reader = pypdf.PdfReader(io.BytesIO(content))
            text_parts = []

            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)

            return "\n\n".join(text_parts)

        except ImportError:
            logger.warning("pypdf not installed, trying PyPDF2")
            try:
                import PyPDF2

                reader = PyPDF2.PdfReader(io.BytesIO(content))
                text_parts = []

                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)

                return "\n\n".join(text_parts)
            except ImportError:
                raise AttachmentProcessorError("No PDF library available")

    def extract_docx(self, content: bytes) -> str:
        """Extract text from DOCX file."""
        try:
            from docx import Document

            doc = Document(io.BytesIO(content))
            text_parts = []

            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        text_parts.append(row_text)

            return "\n\n".join(text_parts)

        except ImportError:
            raise AttachmentProcessorError("python-docx not installed")

    def extract_doc(self, content: bytes) -> str:
        """Extract text from legacy DOC file."""
        # Legacy .doc files are more complex
        # Try using textract or antiword if available
        try:
            import textract

            text = textract.process(
                None,
                extension=".doc",
                input_bytes=content,
            )
            return text.decode("utf-8", errors="replace")
        except ImportError:
            logger.warning("textract not available for .doc files")
            raise UnsupportedFormatError(
                "Legacy .doc files require textract library"
            )

    def extract_xlsx(self, content: bytes) -> str:
        """Extract text from XLSX file."""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            text_parts = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"=== Sheet: {sheet_name} ===")

                for row in sheet.iter_rows():
                    row_values = []
                    for cell in row:
                        if cell.value is not None:
                            row_values.append(str(cell.value))
                    if row_values:
                        text_parts.append(" | ".join(row_values))

            return "\n".join(text_parts)

        except ImportError:
            raise AttachmentProcessorError("openpyxl not installed")

    def extract_xls(self, content: bytes) -> str:
        """Extract text from legacy XLS file."""
        try:
            import xlrd

            wb = xlrd.open_workbook(file_contents=content)
            text_parts = []

            for sheet in wb.sheets():
                text_parts.append(f"=== Sheet: {sheet.name} ===")

                for row_idx in range(sheet.nrows):
                    row_values = []
                    for col_idx in range(sheet.ncols):
                        value = sheet.cell_value(row_idx, col_idx)
                        if value:
                            row_values.append(str(value))
                    if row_values:
                        text_parts.append(" | ".join(row_values))

            return "\n".join(text_parts)

        except ImportError:
            raise AttachmentProcessorError("xlrd not installed")

    def extract_pptx(self, content: bytes) -> str:
        """Extract text from PPTX file."""
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_text.append(f"=== Slide {slide_num} ===")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                    # Extract text from tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join(
                                cell.text.strip() for cell in row.cells
                            )
                            if row_text.strip():
                                slide_text.append(row_text)

                if len(slide_text) > 1:  # More than just the slide header
                    text_parts.extend(slide_text)

            return "\n\n".join(text_parts)

        except ImportError:
            raise AttachmentProcessorError("python-pptx not installed")

    def extract_ppt(self, content: bytes) -> str:
        """Extract text from legacy PPT file."""
        try:
            import textract

            text = textract.process(
                None,
                extension=".ppt",
                input_bytes=content,
            )
            return text.decode("utf-8", errors="replace")
        except ImportError:
            logger.warning("textract not available for .ppt files")
            raise UnsupportedFormatError(
                "Legacy .ppt files require textract library"
            )

    def extract_zip_contents(self, content: bytes) -> str:
        """Extract text from files inside a ZIP archive."""
        import zipfile

        text_parts = []
        try:
            with zipfile.ZipFile(io.BytesIO(content), 'r') as zf:
                for file_info in zf.infolist():
                    if file_info.is_dir():
                        continue

                    filename = file_info.filename
                    ext = Path(filename).suffix.lower()

                    # Only extract text from supported file types
                    if ext in self.EXTENSION_MAP:
                        mime_type = self.EXTENSION_MAP[ext]
                        # Skip nested zips to avoid recursion issues
                        if mime_type == "application/zip":
                            continue

                        try:
                            file_content = zf.read(file_info)
                            processor_name = self.SUPPORTED_TYPES.get(mime_type)
                            if processor_name:
                                processor = getattr(self, processor_name, None)
                                if processor:
                                    extracted = processor(file_content)
                                    if extracted.strip():
                                        text_parts.append(f"=== {filename} ===")
                                        text_parts.append(extracted)
                        except Exception as e:
                            logger.warning(f"Failed to extract {filename} from ZIP: {e}")
                            continue

            return "\n\n".join(text_parts) if text_parts else ""

        except zipfile.BadZipFile:
            raise AttachmentProcessorError("Invalid or corrupted ZIP file")

    def extract_text_content(self, content: bytes) -> str:
        """Extract text from plain text file."""
        # Try different encodings
        encodings = ["utf-8", "latin-1", "cp1252", "iso-8859-1"]

        for encoding in encodings:
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue

        # Fall back to replace errors
        return content.decode("utf-8", errors="replace")

    # Alias for text extraction
    extract_text = extract_text_content

    def extract_html(self, content: bytes) -> str:
        """Extract text from HTML file."""
        try:
            from bs4 import BeautifulSoup

            # Decode content
            text = self.extract_text_content(content)

            # Parse HTML
            soup = BeautifulSoup(text, "html.parser")

            # Remove script and style elements
            for element in soup(["script", "style", "head", "meta", "link"]):
                element.decompose()

            # Get text
            return soup.get_text(separator="\n", strip=True)

        except ImportError:
            # Fall back to basic HTML stripping
            import re

            text = self.extract_text_content(content)
            # Remove HTML tags
            text = re.sub(r"<[^>]+>", " ", text)
            # Clean up whitespace
            text = re.sub(r"\s+", " ", text)
            return text.strip()


# Global instance
attachment_processor = AttachmentProcessor()


def get_attachment_processor() -> AttachmentProcessor:
    """Get the attachment processor instance."""
    return attachment_processor
